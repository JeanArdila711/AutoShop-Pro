"""
build_pptx.py — Genera AutoShop_Pro_Entregable2.pptx
Ejecutar: python docs/build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paleta ──────────────────────────────────────────────────
RED     = RGBColor(0xDC, 0x26, 0x26)
RED_D   = RGBColor(0x99, 0x1B, 0x1B)
DARK    = RGBColor(0x0F, 0x17, 0x2A)
DARK2   = RGBColor(0x1E, 0x29, 0x3B)
SLATE   = RGBColor(0x33, 0x41, 0x55)
GRAY    = RGBColor(0x64, 0x74, 0x8B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE  = RGBColor(0xF9, 0x73, 0x16)
YELLOW  = RGBColor(0xF5, 0x9E, 0x0B)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completamente en blanco

# ────────────────────────────────────────────────────────────
# helpers
# ────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, line_color=None, line_w=Pt(1.5), radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def line_shape(slide, x1, y1, x2, y2, color=GRAY, width=Pt(1.5)):
    """Draw a line using a thin rectangle as workaround (pptx connectors need extra setup)."""
    from pptx.util import Emu
    import math
    # Use add_connector if available, else skip
    try:
        from pptx.enum.shapes import MSO_CONNECTOR
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        connector.line.color.rgb = color
        connector.line.width = width
    except Exception:
        pass

def accent_bar(slide, t, color=RED):
    """Thin horizontal accent bar."""
    box(slide, 0, t, 13.33, 0.06, fill_color=color, line_color=None)

def slide_header(slide, title, subtitle=None, title_color=WHITE, sub_color=None):
    if sub_color is None:
        sub_color = GRAY
    txt(slide, title, 0.6, 0.22, 12, 0.55,
        size=Pt(32), bold=True, color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.6, 0.82, 12, 0.38,
            size=Pt(13), bold=False, color=sub_color, align=PP_ALIGN.LEFT, italic=True)
    accent_bar(slide, 0.72)

def footer(slide, text="AutoShop Pro  ·  Entregable 2  ·  Arquitectura de Software  ·  2026-I"):
    box(slide, 0, 7.2, 13.33, 0.3, fill_color=DARK2)
    txt(slide, text, 0.4, 7.22, 12.5, 0.26,
        size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

def placeholder_img(slide, l, t, w, h, label="📸 Captura aquí", sublabel=""):
    """Dashed placeholder box for screenshots."""
    shape = box(slide, l, t, w, h, fill_color=RGBColor(0x0F,0x17,0x2A),
                line_color=SLATE, line_w=Pt(2))
    # dashed border effect via thin inner box
    box(slide, l+0.05, t+0.05, w-0.1, h-0.1,
        fill_color=None, line_color=RGBColor(0x33,0x41,0x55), line_w=Pt(1))
    txt(slide, label,   l, t + h/2 - 0.3, w, 0.4,
        size=Pt(16), bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    if sublabel:
        txt(slide, sublabel, l, t + h/2 + 0.1, w, 0.3,
            size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)

def tag(slide, l, t, w, h, text, fill, text_color=WHITE):
    b = box(slide, l, t, w, h, fill_color=fill)
    txt(slide, text, l+0.05, t+0.02, w-0.1, h-0.04,
        size=Pt(9), bold=True, color=text_color, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)

# Red accent top
box(s, 0, 0, 13.33, 0.18, fill_color=RED)

# Large title
txt(s, "AutoShop Pro", 0.8, 1.2, 11.7, 1.2,
    size=Pt(60), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Red underline
box(s, 4.5, 2.55, 4.3, 0.08, fill_color=RED)

# Subtitle
txt(s, "Entregable 2 — Evolución a Microservicios en AWS",
    0.8, 2.75, 11.7, 0.6,
    size=Pt(22), bold=False, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)

# Info grid
for i, (label, val) in enumerate([
    ("Curso",    "Arquitectura de Software"),
    ("Profesor", "Nicolás Ramírez Vélez"),
    ("Semestre", "7 · 2026-I"),
    ("Patrón",   "Strangler Fig Pattern"),
]):
    col = 1.6 + i * 2.55
    box(s, col, 3.6, 2.3, 0.9, fill_color=DARK2, line_color=SLATE, line_w=Pt(1))
    txt(s, label, col+0.1, 3.65, 2.1, 0.32,
        size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER, bold=True)
    txt(s, val,   col+0.1, 3.97, 2.1, 0.42,
        size=Pt(11), color=LIGHT, align=PP_ALIGN.CENTER, bold=False)

# Tech badges
for i, (tech, clr) in enumerate([
    ("Django 5.2", BLUE),
    ("Flask", RED),
    ("Redis 7", PURPLE),
    ("Celery", ORANGE),
    ("Nginx", GREEN),
    ("AWS EC2", YELLOW),
    ("Docker", BLUE),
]):
    bx = 0.8 + i * 1.7
    box(s, bx, 4.85, 1.55, 0.42, fill_color=clr)
    txt(s, tech, bx+0.05, 4.88, 1.45, 0.36,
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom info
txt(s, "174.129.161.254  ·  EC2 t3.medium  ·  10 contenedores Docker  ·  us-east-1",
    0.8, 5.55, 11.7, 0.4,
    size=Pt(11), color=GRAY, align=PP_ALIGN.CENTER)

box(s, 0, 7.2, 13.33, 0.3, fill_color=RGBColor(0x09,0x0F,0x18))
txt(s, "EAFIT  ·  Ingeniería de Sistemas  ·  Mayo 2026",
    0.4, 7.22, 12.5, 0.26, size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — ÍNDICE / AGENDA
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "Contenido de la Presentación",
             "Evolución del monolito Django a microservicios en AWS")
footer(s)

items = [
    ("01", "Arquitectura General del Sistema",    "Django + Flask + Redis + Nginx + AWS EC2",           RED,    1.1),
    ("02", "Strangler Fig Pattern",               "Migración progresiva: API v1 (Django) → API v2 (Flask)", BLUE,   2.3),
    ("03", "Comunicación Asíncrona",              "Redis Broker + Celery Worker + eventos entre servicios", PURPLE, 3.5),
    ("04", "Infraestructura AWS",                 "EC2 · VPC · Security Group · 10 contenedores",      GREEN,  4.7),
    ("05", "Evidencia del Sistema en Producción", "Capturas del proyecto funcionando en AWS",            ORANGE, 5.9),
]

for num, title, sub, color, top in items:
    box(s, 0.8, top, 0.7, 0.72, fill_color=color)
    txt(s, num, 0.8, top+0.14, 0.7, 0.44,
        size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.6, top, 10.8, 0.72, fill_color=DARK2, line_color=SLATE, line_w=Pt(1))
    txt(s, title, 1.75, top+0.05, 10.5, 0.36,
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(s, sub,   1.75, top+0.4,  10.5, 0.26,
        size=Pt(10), color=GRAY, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — ARQUITECTURA GENERAL (descripción)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "① Arquitectura General del Sistema", "AWS EC2 · Nginx · Django · 6 Flask microservicios · Redis · Celery")
footer(s)

# Main description box
box(s, 0.5, 1.1, 12.3, 5.9, fill_color=DARK2, line_color=SLATE, line_w=Pt(1))

# Columns
# Left: layers
layers = [
    ("🌐 Internet / Clientes",      "Browser, API consumers, Equipo aliado",           GRAY,   DARK),
    ("⚡ Nginx 1.25 (API Gateway)", "Strangler Pattern: /api/v1/* → Django · /api/v2/* → Flask", GREEN, DARK),
    ("🐍 Django 5.2 (Monolito)",    "Frontend HTML · API v1 · i18n ES/EN · Adapter ExchangeRate", BLUE,  DARK),
    ("🔴 Redis 7 (Broker)",         "Task Queue (Celery) · Pub/Sub eventos entre servicios",       PURPLE,DARK),
    ("⚙ Celery Worker",            "Async: PDF facturación · Email notificaciones · SendGrid",    ORANGE,DARK),
]
for i, (title, desc, color, bg_c) in enumerate(layers):
    top = 1.25 + i * 1.0
    box(s, 0.65, top, 0.22, 0.72, fill_color=color)
    txt(s, title, 0.95, top+0.04, 5.2, 0.34, size=Pt(11), bold=True, color=WHITE)
    txt(s, desc,  0.95, top+0.38, 5.2, 0.28, size=Pt(9),  color=GRAY)
    if i < len(layers)-1:
        txt(s, "↓", 0.7, top+0.76, 0.22, 0.22, size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)

# Right: microservices grid
micros = [
    ("ordenes",       ":5001", "WorkOrder · Bahia"),
    ("inventario",    ":5002", "Partes · OrdenCompra"),
    ("facturacion",   ":5003", "Factura · PDF async"),
    ("citas",         ":5004", "Agenda · State Machine"),
    ("notificaciones",":5005", "Email · SendGrid"),
    ("predictivo",    ":5000", "ML scoring · Alertas"),
]
txt(s, "Flask Microservicios (API v2)", 6.8, 1.18, 5.8, 0.36,
    size=Pt(12), bold=True, color=RED, align=PP_ALIGN.CENTER)

for i, (name, port, desc) in enumerate(micros):
    col = 6.8 + (i % 2) * 2.95
    row = 1.65 + (i // 2) * 1.55
    box(s, col, row, 2.75, 1.35, fill_color=DARK, line_color=RED, line_w=Pt(1.5))
    txt(s, name,  col+0.12, row+0.1,  2.5, 0.36, size=Pt(11), bold=True, color=RGBColor(0xF8,0x71,0x71))
    txt(s, port,  col+0.12, row+0.46, 2.5, 0.26, size=Pt(9),  color=GRAY)
    txt(s, desc,  col+0.12, row+0.72, 2.5, 0.26, size=Pt(9),  color=LIGHT)
    txt(s, "✅ healthy", col+0.12, row+1.02, 2.5, 0.24, size=Pt(9), color=GREEN)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — DIAGRAMA ARQUITECTURA GENERAL (imagen placeholder)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "① Arquitectura General — Diagrama",
             "Ver: docs/diagramas.html → Diagrama ①")
footer(s)

placeholder_img(s, 0.5, 1.1, 12.3, 6.0,
                "📸  Insertar captura del Diagrama ①",
                "docs/diagramas.html → sección 'Arquitectura General'")

# ════════════════════════════════════════════════════════════
# SLIDE 5 — STRANGLER PATTERN (descripción)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "② Strangler Fig Pattern",
             "Migración progresiva del monolito a microservicios sin downtime")
footer(s)

# Concept box
box(s, 0.5, 1.1, 12.3, 1.15, fill_color=RGBColor(0x1C,0x17,0x00), line_color=YELLOW, line_w=Pt(1.5))
txt(s, "¿Qué es el Strangler Fig Pattern?",
    0.7, 1.15, 11.9, 0.38, size=Pt(13), bold=True, color=YELLOW)
txt(s, "Nginx intercepta TODAS las peticiones. Las rutas /api/v1/* siguen al monolito Django (legacy). "
       "Las rutas /api/v2/<servicio>/* se redirigen al microservicio Flask correspondiente. "
       "El monolito 'se estrangula' gradualmente sin interrumpir el servicio.",
    0.7, 1.53, 11.9, 0.65, size=Pt(10), color=LIGHT)

# Three columns: Legacy | Nginx | Micros
for col_x, col_title, col_color, items_list in [
    (0.5,  "MONOLITO (Legacy)",    BLUE,   [
        "/workorders/       → DashboardView",
        "/vehicles/         → VehicleListView",
        "/owners/           → OwnerListView",
        "/mechanics/        → MechanicListView",
        "/api/v1/workorders → API REST Django",
        "/api/v1/vehicles   → API REST Django",
        "HTML Templates + i18n ES/EN",
        "ExchangeRate Adapter Pattern",
    ]),
    (4.55, "NGINX — INTERCEPTOR",  YELLOW, [
        "location /api/v1/ {",
        "  proxy_pass django:8000",
        "}",
        "location /api/v2/ordenes/ {",
        "  proxy_pass ordenes:5001",
        "}",
        "location /api/v2/inventario/ {",
        "  proxy_pass inventario:5002",
        "}  ... (6 reglas)",
    ]),
    (8.6,  "MICROSERVICIOS (Nuevo)", RED,  [
        "/api/v2/ordenes/*       :5001",
        "/api/v2/inventario/*    :5002",
        "/api/v2/facturacion/*   :5003",
        "/api/v2/citas/*         :5004",
        "/api/v2/notificaciones/ :5005",
        "/api/v2/predictivo/*    :5000",
        "/api/v2/catalogo/       (aliada)",
        "Migración: 70% ✅",
    ]),
]:
    box(s, col_x, 2.4, 4.0, 4.7, fill_color=DARK2, line_color=col_color, line_w=Pt(1.5))
    txt(s, col_title, col_x+0.15, 2.45, 3.7, 0.38,
        size=Pt(11), bold=True, color=col_color)
    for j, item in enumerate(items_list):
        txt(s, item, col_x+0.2, 2.88 + j*0.5, 3.6, 0.44,
            size=Pt(9), color=LIGHT if "}" not in item else GRAY)

# Progress
box(s, 0.5, 7.0, 12.3, 0.12, fill_color=DARK2)
box(s, 0.5, 7.0, 8.61, 0.12, fill_color=GREEN)
txt(s, "Progreso de migración: 70%  (6 microservicios migrados, frontend permanece en Django)",
    0.5, 6.8, 12.3, 0.2, size=Pt(9), color=GREEN)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — STRANGLER PATTERN (diagrama placeholder)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "② Strangler Fig Pattern — Diagrama",
             "Ver: docs/diagramas.html → Diagrama ②")
footer(s)

placeholder_img(s, 0.5, 1.1, 12.3, 6.0,
                "📸  Insertar captura del Diagrama ②",
                "docs/diagramas.html → sección 'Strangler Fig Pattern'")

# ════════════════════════════════════════════════════════════
# SLIDE 7 — COMUNICACIÓN ASÍNCRONA (descripción)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "③ Comunicación Asíncrona",
             "Redis Broker + Celery Worker — eventos entre productores y consumidores")
footer(s)

# Concept
box(s, 0.5, 1.1, 12.3, 0.72, fill_color=RGBColor(0x15,0x0A,0x2E), line_color=PURPLE, line_w=Pt(1.5))
txt(s, "Patrón: Productor → Cola Redis → Celery Worker → Servicio externo",
    0.7, 1.14, 11.9, 0.32, size=Pt(12), bold=True, color=PURPLE)
txt(s, "Las tareas pesadas (PDF, emails, notificaciones) se procesan de forma asíncrona sin bloquear el request HTTP.",
    0.7, 1.46, 11.9, 0.3, size=Pt(10), color=LIGHT)

# Three columns
cols = [
    (0.5,  "PRODUCTORES",  BLUE,   [
        "Django",         "orden.creada → redis",
        "facturacion",    "factura.generada → redis",
        "citas",          "cita.confirmada → redis",
        "ordenes",        "orden.completada → redis",
    ]),
    (4.55, "REDIS (Broker)", PURPLE, [
        "celery (queue)",      "Task Queue principal",
        "celery_results",      "Resultados de tareas",
        "ordenes.eventos",     "Pub/Sub canal",
        "facturacion.eventos", "Pub/Sub canal",
    ]),
    (8.6,  "CONSUMIDORES",  ORANGE, [
        "Celery Worker",       "4 workers concurrentes",
        "enviar_email()",      "@app.task",
        "generar_pdf()",       "@app.task",
        "notificar_prop()",    "@app.task",
    ]),
]

for col_x, col_title, col_color, items_list in cols:
    box(s, col_x, 1.95, 4.0, 4.9, fill_color=DARK2, line_color=col_color, line_w=Pt(1.5))
    txt(s, col_title, col_x+0.15, 2.0, 3.7, 0.38, size=Pt(11), bold=True, color=col_color)
    for j in range(0, len(items_list), 2):
        name = items_list[j]
        desc = items_list[j+1] if j+1 < len(items_list) else ""
        row_top = 2.46 + (j//2) * 1.05
        box(s, col_x+0.15, row_top, 3.7, 0.9, fill_color=DARK, line_color=SLATE, line_w=Pt(1))
        txt(s, name, col_x+0.28, row_top+0.08, 3.4, 0.36, size=Pt(11), bold=True, color=WHITE)
        txt(s, desc, col_x+0.28, row_top+0.5,  3.4, 0.28, size=Pt(9),  color=GRAY)

# Arrows between columns
txt(s, "→", 4.35, 3.8, 0.4, 0.5, size=Pt(28), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
txt(s, "→", 8.4,  3.8, 0.4, 0.5, size=Pt(28), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# External service
box(s, 0.5, 7.0, 12.3, 0.38, fill_color=RGBColor(0x1C,0x17,0x00), line_color=YELLOW, line_w=Pt(1))
txt(s, "Servicio externo final: SendGrid (email) · Twilio (SMS futuro) · S3 (PDF storage futuro)",
    0.7, 7.04, 11.9, 0.28, size=Pt(10), color=YELLOW)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — COMUNICACIÓN ASÍNCRONA (diagrama)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "③ Comunicación Asíncrona — Diagrama",
             "Ver: docs/diagramas.html → Diagrama ③")
footer(s)

placeholder_img(s, 0.5, 1.1, 12.3, 6.0,
                "📸  Insertar captura del Diagrama ③",
                "docs/diagramas.html → sección 'Comunicación Asíncrona'")

# ════════════════════════════════════════════════════════════
# SLIDE 9 — INFRAESTRUCTURA AWS (descripción)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "④ Infraestructura en AWS Academy",
             "EC2 t3.medium · us-east-1 · 10 contenedores Docker · IP: 174.129.161.254")
footer(s)

# AWS stack visual
layers_aws = [
    ("☁  AWS Academy Cloud",                   "Account: 975049893369",                      YELLOW, RGBColor(0x1C,0x17,0x00)),
    ("🌎  Region: us-east-1 (N. Virginia)",     "Capa de disponibilidad geográfica",          ORANGE, RGBColor(0x1A,0x0D,0x00)),
    ("🔒  VPC: vpc-0fe167703ba7486d8",          "CIDR: 172.31.0.0/16  ·  Red privada virtual", BLUE,  DARK2),
    ("📦  Subnet: subnet-0b0086b61e9c991e8",    "AZ: us-east-1a  ·  Subred pública",          BLUE,  DARK2),
    ("🛡  Security Group: autoshop-pro-sg",     "Inbound: :22 SSH · :80 HTTP · :8000-8006",    RED,   DARK2),
    ("🖥  EC2: i-0eabd783d5f76efe2",           "t3.medium · 2vCPU · 4GB RAM · 20GB gp3 · Amazon Linux 2023", GREEN, RGBColor(0x0A,0x1A,0x0A)),
    ("🐳  Docker Engine 25.0 + Compose v5",    "10 contenedores · autoshop_net (bridge) · Todos healthy ✅", RGBColor(0x06,0xB6,0xD4), DARK2),
]

for i, (title, desc, color, bg_c) in enumerate(layers_aws):
    indent = i * 0.15
    top = 1.18 + i * 0.74
    w = 12.3 - indent * 2
    box(s, 0.5 + indent, top, w, 0.64, fill_color=bg_c, line_color=color, line_w=Pt(1.5))
    txt(s, title, 0.68 + indent, top+0.06, w * 0.45, 0.3,  size=Pt(10), bold=True, color=color)
    txt(s, desc,  0.68 + indent + w*0.45, top+0.06, w*0.53, 0.3, size=Pt(9), color=GRAY)

# Containers summary
box(s, 0.5, 6.5, 12.3, 0.6, fill_color=DARK2, line_color=GREEN, line_w=Pt(1))
containers = ["nginx:80","django:8000","redis:6379","celery","ordenes:5001",
              "inventario:5002","facturacion:5003","citas:5004","notificaciones:5005","predictivo:5000"]
txt(s, "Contenedores: " + "  ·  ".join(containers),
    0.7, 6.56, 11.9, 0.42, size=Pt(8.5), color=GREEN)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — INFRAESTRUCTURA AWS (diagrama)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "④ Infraestructura AWS — Diagrama",
             "Ver: docs/diagramas.html → Diagrama ④")
footer(s)

placeholder_img(s, 0.5, 1.1, 12.3, 6.0,
                "📸  Insertar captura del Diagrama ④",
                "docs/diagramas.html → sección 'Infraestructura AWS'")

# ════════════════════════════════════════════════════════════
# SLIDE 11 — EVIDENCIA: Dashboard / Home
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "⑤ Evidencia — Dashboard en Producción",
             "http://174.129.161.254/workorders/  ·  AWS EC2 us-east-1")
footer(s)

placeholder_img(s, 0.5, 1.1, 12.3, 5.5,
                "📸  Captura: Dashboard AutoShop Pro en AWS",
                "http://174.129.161.254/workorders/")

box(s, 0.5, 6.72, 12.3, 0.42, fill_color=DARK2, line_color=SLATE, line_w=Pt(1))
txt(s, "URL pública: http://174.129.161.254/workorders/  ·  Django 5.2  ·  Nginx gateway  ·  AWS EC2 t3.medium",
    0.7, 6.77, 11.9, 0.3, size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 12 — EVIDENCIA: Microservicios health
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "⑤ Evidencia — Health Endpoints Microservicios",
             "Todos los servicios responden HTTP 200 en producción")
footer(s)

# Left: screenshot placeholder
placeholder_img(s, 0.5, 1.1, 6.0, 5.5,
                "📸  Captura: /api/v2/*/health",
                "Respuestas JSON de los 6 microservicios")

# Right: endpoint list
box(s, 6.7, 1.1, 6.1, 5.5, fill_color=DARK2, line_color=SLATE, line_w=Pt(1))
txt(s, "Endpoints verificados ✅", 6.9, 1.18, 5.7, 0.38,
    size=Pt(12), bold=True, color=GREEN)

endpoints = [
    ("/api/v2/inventario/health",    "HTTP 200", "inventario_service :5002"),
    ("/api/v2/citas/health",         "HTTP 200", "citas_service :5004"),
    ("/api/v2/ordenes/health",       "HTTP 200", "ordenes_service :5001"),
    ("/api/v2/facturacion/health",   "HTTP 200", "facturacion_service :5003"),
    ("/api/v2/notificaciones/health","HTTP 200", "notificaciones_service :5005"),
    ("/api/v2/predictivo/health",    "HTTP 200", "predictivo_service :5000"),
    ("/api/v2/catalogo/",            "HTTP 200", "API pública equipo aliado"),
]
for i, (ep, status, svc) in enumerate(endpoints):
    top = 1.65 + i * 0.68
    box(s, 6.8, top, 5.8, 0.58, fill_color=DARK, line_color=RGBColor(0x16,0xA3,0x4A), line_w=Pt(1))
    txt(s, ep,     6.95, top+0.04, 5.5, 0.26, size=Pt(9),  bold=True, color=WHITE)
    txt(s, status, 6.95, top+0.3,  1.5, 0.2,  size=Pt(9),  color=GREEN, bold=True)
    txt(s, svc,    8.55, top+0.3,  4.0, 0.2,  size=Pt(8),  color=GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 13 — EVIDENCIA: API Catálogo + i18n
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "⑤ Evidencia — API Catálogo & i18n ES/EN",
             "Endpoint público para equipo aliado · Soporte bilingüe Django")
footer(s)

# Two placeholders side by side
placeholder_img(s, 0.5, 1.1, 6.0, 5.5,
                "📸  Captura: GET /api/v2/catalogo/",
                "JSON con partes mecánicas para equipo aliado")

placeholder_img(s, 6.7, 1.1, 6.1, 5.5,
                "📸  Captura: Dashboard en inglés (EN)",
                "Selector de idioma ES/EN funcionando")

# ════════════════════════════════════════════════════════════
# SLIDE 14 — EVIDENCIA: Docker containers
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
slide_header(s, "⑤ Evidencia — Contenedores Docker en AWS",
             "docker compose ps — 10 contenedores, todos healthy")
footer(s)

placeholder_img(s, 0.5, 1.1, 12.3, 4.8,
                "📸  Captura: docker compose ps (SSH en EC2)",
                "Muestra los 10 contenedores con estado 'healthy'")

# Container status grid
containers_info = [
    ("autoshop_nginx",          "Up · :80",  "healthy"),
    ("autoshop_django",         "Up · :8000","healthy"),
    ("autoshop_redis",          "Up · :6379","healthy"),
    ("autoshop_celery",         "Up",        "healthy"),
    ("autoshop_ordenes",        "Up · :5001","healthy"),
    ("autoshop_inventario",     "Up · :5002","healthy"),
    ("autoshop_facturacion",    "Up · :5003","healthy"),
    ("autoshop_citas",          "Up · :5004","healthy"),
    ("autoshop_notificaciones", "Up · :5005","healthy"),
    ("autoshop_predictivo",     "Up · :5000","healthy"),
]
for i, (name, status, health) in enumerate(containers_info):
    col = 0.5 + (i % 5) * 2.45
    row = 6.1 + (i // 5) * 0.52
    box(s, col, row, 2.35, 0.44, fill_color=DARK2, line_color=GREEN, line_w=Pt(1))
    txt(s, name,   col+0.1, row+0.03, 2.1, 0.22, size=Pt(7.5), bold=True, color=WHITE)
    txt(s, status, col+0.1, row+0.25, 1.2, 0.16, size=Pt(7),   color=GRAY)
    txt(s, "✅",   col+1.9, row+0.1,  0.4, 0.28, size=Pt(12),  color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 15 — CIERRE / RESUMEN
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
box(s, 0, 0, 13.33, 0.18, fill_color=RED)
footer(s)

txt(s, "Resumen del Entregable 2", 0.8, 0.45, 11.7, 0.72,
    size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 4.5, 1.22, 4.3, 0.08, fill_color=RED)

achievements = [
    ("✅", "Strangler Fig Pattern",      "Nginx enruta v1→Django / v2→Flask sin downtime",          RED),
    ("✅", "6 Microservicios Flask",     "ordenes · inventario · facturación · citas · notificaciones · predictivo", BLUE),
    ("✅", "Comunicación Asíncrona",     "Redis Broker + Celery Worker · PDF · Email · eventos",     PURPLE),
    ("✅", "Adapter Pattern",            "ExchangeRate-API (real/mock) + SendGrid (email/mock)",     ORANGE),
    ("✅", "i18n Bilingüe ES/EN",        "Django gettext · LocaleMiddleware · .po/.mo compilados",   GREEN),
    ("✅", "Deploy AWS EC2",             "t3.medium · us-east-1 · 174.129.161.254 · 10 contenedores",YELLOW),
    ("✅", "API Pública Equipo Aliado",  "GET /api/v2/catalogo/ · documentada y desplegada en AWS",  RGBColor(0x06,0xB6,0xD4)),
]

for i, (icon, title, desc, color) in enumerate(achievements):
    col = 0.5 + (i % 2) * 6.35
    row = 1.55 + (i // 2) * 1.3
    if i == 6:  # último centrado
        col = 3.42
    box(s, col, row, 6.1, 1.12, fill_color=DARK2, line_color=color, line_w=Pt(1.5))
    txt(s, icon,  col+0.15, row+0.2, 0.6, 0.7, size=Pt(20), align=PP_ALIGN.CENTER)
    txt(s, title, col+0.8,  row+0.08, 5.1, 0.36, size=Pt(11), bold=True, color=color)
    txt(s, desc,  col+0.8,  row+0.5,  5.1, 0.5,  size=Pt(9),  color=GRAY)

txt(s, "http://174.129.161.254/",
    0.8, 7.05, 11.7, 0.3,
    size=Pt(12), bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "AutoShop_Pro_Entregable2.pptx")
prs.save(out)
print(f"OK  Guardado: {out}")
print(f"    Diapositivas: {len(prs.slides)}")
